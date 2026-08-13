"""Skill 注册表：工具定义 + 本地执行。

所有文件操作都被限制在 config.WORKSPACE_ROOT 沙箱内，防止越权读写。
OpenAI 兼容接口的 tool 定义可让模型自行决定是否调用。
"""
from __future__ import annotations

import asyncio
import contextvars
import hashlib
import json
import os
import re
import subprocess
import signal
import logging
from pathlib import Path
from typing import Any, Callable

import httpx

from app.config import config
from app import db

# ---- 会话级「工作目录范围」覆盖 ----
# 某次对话若设定了 work_dir，则该对话内的文件工具以它为根，
# 而非全局的 WORKSPACE_ROOT / TARGET_DIR。用 ContextVar 承载：
# 在 chat_stream 的协程里 set，asyncio.to_thread 会复制上下文，工具线程内可读到。
_WORK_DIR: contextvars.ContextVar[str | None] = contextvars.ContextVar(
    "mb_work_dir", default=None
)


def set_work_dir(path: str | None) -> None:
    _WORK_DIR.set(path)


def get_work_dir() -> str | None:
    return _WORK_DIR.get()


def effective_root() -> Path:
    """文件读写/列目录/搜索/解析文档的有效根（work_dir 优先，否则全局 WORKSPACE_ROOT）。"""
    wd = _WORK_DIR.get()
    return Path(wd).resolve() if wd else config.WORKSPACE_ROOT


def effective_target_dir() -> Path:
    """manage_dir 的有效根（work_dir 优先，否则全局 TARGET_DIR）。"""
    wd = _WORK_DIR.get()
    return Path(wd).resolve() if wd else config.TARGET_DIR

# ---- 安全审计与命令沙箱 ----
_audit = logging.getLogger("mbuddy_audit")
_audit_inited = False
def _init_audit():
    global _audit_inited
    if _audit_inited:
        return
    _audit_inited = True
    try:
        log_path = str(config.DB_PATH).rsplit(".", 1)[0] + ".audit.log"
        _h = logging.FileHandler(log_path, encoding="utf-8")
        _h.setFormatter(logging.Formatter("%(asctime)s %(message)s"))
        _audit.addHandler(_h)
        _audit.setLevel(logging.INFO)
    except Exception:  # noqa: BLE001
        pass
def audit(action: str, detail: str) -> None:
    try:
        _init_audit()
        _audit.info("%s | %s", action, detail)
    except Exception:  # noqa: BLE001
        pass

# 危险命令基名黑名单（命中即拒绝，防提示注入启动新 shell / 下载执行 / 反弹等）
_DANGEROUS_BIN = {
    "powershell", "pwsh", "cmd", "bash", "sh", "certutil", "curl", "wget",
    "nc", "ncat", "netcat", "telnet", "reg", "schtasks", "at", "sc",
    "tscon", "psexec", "ftp",
}
# 命令链分隔符（裸用易被注入拼接额外命令）；保留 | 与 >（常用且危害较小）
_CHAIN_RE = re.compile(r";|&&|\|\||(?<!&)&(?!&)")

def _run_subprocess(cmd, cwd, timeout: int = 30, shell: bool = False) -> str:
    """统一执行子进程：整组清理（超时杀进程树）+ 输出截断，防孤儿进程与海量输出。"""
    try:
        proc = subprocess.run(
            cmd, shell=shell, cwd=str(cwd),
            capture_output=True, text=True,
            timeout=timeout, start_new_session=(os.name != "nt"),
        )
        out = (proc.stdout or "") + (proc.stderr or "")
        return out[:65536].strip() or "[无输出]"
    except subprocess.TimeoutExpired:
        try:
            if os.name == "nt":
                subprocess.run(
                    ["taskkill", "/F", "/T", "/PID", str(proc.pid)],
                    capture_output=True,
                )
            else:
                os.killpg(os.getpgid(proc.pid), signal.SIGKILL)
        except Exception:  # noqa: BLE001
            try:
                proc.kill()
            except Exception:  # noqa: BLE001
                pass
        return f"[错误] 执行超时({timeout}s)，已终止进程组"
    except Exception as e:  # noqa: BLE001
        return f"[错误] 执行异常: {e}"


# 允许执行的命令白名单（用于 run_command / run_code 的底层调用）
_CODE_EXT = {".py": ["python"], ".js": ["node"], ".ts": ["npx", "ts-node"]}


def _safe_path(path: str) -> Path:
    """解析路径。沙箱开启时限制在工作区内；桌面版沙箱关闭时允许任意本地绝对路径。

    相对路径的基准为 effective_root()（若对话设定了 work_dir，则以其为根）。
    """
    root = effective_root()
    raw = Path(path)
    if raw.is_absolute():
        p = raw.resolve()
        if config.MB_SANDBOX and root not in p.parents and p != root:
            raise ValueError(f"路径越界，禁止访问沙箱外: {path}")
        return p
    p = (root / path).resolve()
    if config.MB_SANDBOX and root not in p.parents and p != root:
        try:
            p.relative_to(root)
        except ValueError:
            raise ValueError(f"路径越界，禁止访问沙箱外: {path}")
    return p


def _display(p: Path) -> str:
    """沙箱开启时用相对路径展示，关闭时用绝对路径（桌面版选的真实文件）。"""
    if config.MB_SANDBOX:
        try:
            return str(p.relative_to(effective_root()))
        except ValueError:
            return str(p)
    return str(p)


def read_file(path: str) -> str:
    p = _safe_path(path)
    if not p.is_file():
        return f"[错误] 文件不存在: {path}"
    # 桌面版沙箱关闭时允许任意绝对路径，记录审计日志以防模型越权读取敏感文件
    if not config.MB_SANDBOX:
        audit("read_file", str(p))
    try:
        return p.read_text(encoding="utf-8", errors="replace")
    except Exception as e:  # noqa: BLE001
        return f"[错误] 读取失败: {e}"


def write_file(path: str, content: str) -> str:
    p = _safe_path(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(content, encoding="utf-8")
    return f"已写入 {len(content)} 字符 -> {_display(p)}"


def list_dir(path: str = ".") -> str:
    p = _safe_path(path)
    if not p.is_dir():
        return f"[错误] 目录不存在: {path}"
    entries = []
    for child in sorted(p.iterdir()):
        kind = "目录" if child.is_dir() else "文件"
        entries.append(f"{kind:4} {child.name}")
    return "\n".join(entries) if entries else "(空目录)"


def search_text(pattern: str, path: str = ".") -> str:
    """简单子串搜索（大小写不敏感），返回匹配文件及行号。"""
    root = _safe_path(path)
    if not root.exists():
        return f"[错误] 路径不存在: {path}"
    needle = pattern.lower()
    results = []
    for fp in root.rglob("*"):
        if not fp.is_file():
            continue
        if fp.suffix in {".pyc", ".db", ".sqlite", ".wasm"}:
            continue
        try:
            text = fp.read_text(encoding="utf-8", errors="ignore")
        except Exception:  # noqa: BLE001
            continue
        for i, line in enumerate(text.splitlines(), 1):
            if needle in line.lower():
                rel = _display(fp)
                results.append(f"{rel}:{i}: {line.strip()[:200]}")
                if len(results) >= 50:
                    return "\n".join(results) + "\n[截断] 仅显示前 50 条"
    return "\n".join(results) if results else f"未找到包含 '{pattern}' 的内容"


def run_command(command: str) -> str:
    """执行受控 shell 命令（在工作区内），返回 stdout/stderr。

    含 Windows 命令护栏：在本机为 Windows 时，拦截常见的 Unix-only 命令，
    避免模型误调 ls/cat/rm 等导致失败；并对高危删除命令做确认拦截。
    """
    if config.IS_WINDOWS:
        _UNIX_ONLY = {
            "ls": "dir",
            "cat": "type",
            "rm": "del",
            "rmdir": "rmdir /s /q",
            "cp": "copy",
            "mv": "move",
            "pwd": "cd",
            "grep": "findstr",
            "touch": None,  # Windows 无对应，需用重定向创建
            "clear": "cls",
        }
        first = command.strip().split()[0] if command.strip() else ""
        base = first.split(".")[0].lower()
        if base in _UNIX_ONLY:
            repl = _UNIX_ONLY[base]
            if repl is None:
                return (
                    f"[已拦截] 命令 '{first}' 是 Unix 专用命令，本机为 Windows 不支持。"
                    "如需创建空文件，请使用 write_file 工具；其他操作请改用 Windows 等效命令。"
                )
            return (
                f"[已拦截并提示] 本机为 Windows，'{first}' 是 Unix 命令。请改用 Windows 等效命令: '{repl}'。"
                "（建议优先使用 list_dir/read_file/write_file/manage_dir 等专用工具，而非 shell 命令）"
            )
        # 危险命令基名黑名单（防提示注入启动新 shell / 下载执行 / 反弹等）
        base = first.split(".")[0].lower().replace('"', "").replace("'", "")
        if "\\" in base or "/" in base:
            base = base.rsplit("\\", 1)[-1].rsplit("/", 1)[-1]
        if base in _DANGEROUS_BIN:
            return (
                f"[已拦截] 禁止执行危险命令 '{first}'，以防提示注入执行任意指令。"
                "如需运行脚本，请使用 run_code 工具。"
            )
        # 命令链分隔符禁用（防注入拼接额外命令）
        if _CHAIN_RE.search(command):
            return (
                "[已拦截] 命令包含链式操作符（; && || &），已禁止以防止命令注入。"
                "如需组合操作，请拆成多条命令或改用专用工具。"
            )
        # 高危删除类命令二次确认拦截
        if re.search(r"\b(del|rmdir|rd)\b", command, re.I) and (
            "*" in command or "/" in command or ".." in command
        ):
            return (
                "[已拦截] 检测到高危通配/越级删除命令，已阻止执行以保障本机文件安全。"
                "如需删除，请使用 manage_dir(action='delete') 工具，限定在工作目录范围内。"
            )

    return _run_subprocess(command, effective_root(), timeout=30, shell=True)


def run_code(code: str, lang: str = "python") -> str:
    """将代码写入临时文件并执行，返回结果。"""
    ext = ".py" if lang.lower().startswith("py") else ".js"
    target = _safe_path(f".wb_runs/tmp{ext}")
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_text(code, encoding="utf-8")
    runner = _CODE_EXT.get(ext, ["python"])[0]
    exe = "python" if runner == "python" else runner
    if exe == "python":
        exe = os.environ.get("PYTHON_EXE", "python")
    return _run_subprocess([exe, str(target)], effective_root(), timeout=30, shell=False)


# ---- 新增 Skill：联网检索 ----
async def _web_search_async(query: str, max_results: int = 5) -> str:
    """使用 DuckDuckGo HTML 接口做轻量联网检索（无需 API key）。"""
    try:
        async with httpx.AsyncClient(timeout=15, follow_redirects=True) as client:
            headers = {"User-Agent": "Mozilla/5.0 (compatible; mini-workbuddy/0.1)"}
            resp = await client.post(
                "https://html.duckduckgo.com/html/",
                data={"q": query},
                headers=headers,
            )
            resp.raise_for_status()
            html = resp.text
    except Exception as e:  # noqa: BLE001
        return f"[错误] 联网检索失败: {e}"

    # 解析结果条目
    results = []
    for m in re.finditer(
        r'class="result__a"[^>]*href="([^"]+)"[^>]*>(.*?)</a>.*?class="result__snippet"[^>]*>(.*?)</a>',
        html,
        re.S,
    ):
        url = m.group(1)
        title = re.sub(r"<[^>]+>", "", m.group(2)).strip()
        snippet = re.sub(r"<[^>]+>", "", m.group(3)).strip()
        results.append(f"- {title}\n  {snippet}\n  {url}")
        if len(results) >= max_results:
            break
    return "\n".join(results) if results else "未检索到相关结果。"


async def web_search(query: str, max_results: int = 5) -> str:
    """联网检索公开信息（无需 API key，使用 DuckDuckGo）。

    改为异步实现：原同步包装里调用 asyncio.run() 会在已有事件循环
    （FastAPI/uvicorn）中直接抛 RuntimeError，导致联网检索必挂。
    """
    try:
        return await _web_search_async(query, int(max_results))
    except Exception as e:  # noqa: BLE001
        return (
            f"[错误] 联网检索失败: {e}（DuckDuckGo 接口可能限流或改版，"
            "请稍后重试、换关键词，或改用本地 RAG 检索文档内容。）"
        )


# ---- 新增 Skill：操作特定目录 ----
def _safe_target_path(path: str) -> Path:
    """将相对于有效 TARGET_DIR 的路径解析出来，越界则拒绝（独立于工作区沙箱）。

    若对话设定了 work_dir，有效 TARGET_DIR 即为该目录。
    """
    base = effective_target_dir()
    p = (base / path).resolve()
    try:
        p.relative_to(base)
    except ValueError:
        raise ValueError(f"越界：仅允许操作工作目录范围({base})内的路径: {path}")
    return p


def manage_dir(action: str, path: str, content: str = "") -> str:
    """对 TARGET_DIR 内的目录/文件进行受限操作：list / read / write / mkdir / delete / move。

    action: list | read | write | mkdir | delete | move
    path: 相对 TARGET_DIR 的路径
    content: write 时必填
    dest: move 时的目标相对路径
    """
    try:
        target = _safe_target_path(path)
    except ValueError as e:
        return f"[错误] {e}"

    action = action.lower()
    if action == "list":
        if not target.exists():
            return f"[错误] 路径不存在: {path}"
        if target.is_file():
            return f"文件: {path} ({target.stat().st_size} 字节)"
        entries = [
            f"{'目录' if c.is_dir() else '文件':4} {c.name}"
            for c in sorted(target.iterdir())
        ]
        return "\n".join(entries) if entries else "(空)"
    if action == "read":
        if not target.is_file():
            return f"[错误] 文件不存在: {path}"
        return target.read_text(encoding="utf-8", errors="replace")
    if action == "write":
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(content, encoding="utf-8")
        return f"已写入 {len(content)} 字符 -> {path}"
    if action == "mkdir":
        target.mkdir(parents=True, exist_ok=True)
        return f"已创建目录 -> {path}"
    if action == "delete":
        if not target.exists():
            return f"[错误] 路径不存在: {path}"
        if target.is_dir():
            import shutil

            shutil.rmtree(target)
        else:
            target.unlink()
        return f"已删除 -> {path}"
    if action == "move":
        dest = _safe_target_path(content) if content else None
        if not dest:
            return "[错误] move 需提供 dest（目标相对路径填在 content）"
        target.replace(dest)
        return f"已移动 {path} -> {dest.relative_to(effective_target_dir())}"
    return f"[错误] 未知 action: {action}"


# ---- 新增 Skill：RAG 检索本地文档 ----
_RAG_INDEX_FILE = config.DOC_ROOT / ".rag_index.json"
_RAG_CHUNK = 600
_RAG_OVERLAP = 80


def _chunk_text(text: str, size: int = _RAG_CHUNK, overlap: int = _RAG_OVERLAP) -> list[str]:
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []
    if len(text) <= size:
        return [text]
    chunks, start = [], 0
    while start < len(text):
        chunks.append(text[start : start + size])
        start += size - overlap
    return chunks


def _simple_tokenize(text: str) -> set[str]:
    # 中文按字符、英文按词，做轻量归一化
    text = text.lower()
    toks = set(re.findall(r"[a-z0-9]+", text))
    for ch in text:
        if "\u4e00" <= ch <= "\u9fff":
            toks.add(ch)
    return toks


def rag_index() -> str:
    """扫描 DOC_ROOT 下的文档（md/txt/py/json），分词建立倒排索引并持久化。"""
    config.DOC_ROOT.mkdir(parents=True, exist_ok=True)
    index: dict[str, list[dict]] = {}
    _TEXT_SUF = {".md", ".txt", ".py", ".json", ".csv", ".log", ".yaml", ".yml", ".toml"}
    _DOC_SUF = {".docx", ".xlsx", ".pptx", ".pdf"}
    files = [
        p
        for p in config.DOC_ROOT.rglob("*")
        if p.is_file() and p.suffix.lower() in _TEXT_SUF | _DOC_SUF
    ]
    count = 0
    for fp in files:
        try:
            if fp.suffix.lower() in _DOC_SUF:
                text = parse_document(str(fp), max_chars=50000)
                if text.startswith(("[错误]", "[不支持]", "[内容为空]")):
                    continue
            else:
                text = fp.read_text(encoding="utf-8", errors="ignore")
        except Exception:  # noqa: BLE001
            continue
        rel = str(fp.relative_to(config.DOC_ROOT))
        for i, chunk in enumerate(_chunk_text(text)):
            cid = f"{rel}#{i}"
            for tok in _simple_tokenize(chunk):
                index.setdefault(tok, []).append({"cid": cid, "chunk": chunk[:400]})
            count += 1
    _RAG_INDEX_FILE.write_text(
        json.dumps(index, ensure_ascii=False), encoding="utf-8"
    )
    return f"索引完成：{len(files)} 个文件，{count} 个分块。"


def rag_query(query: str, top_k: int = 5) -> str:
    """在已建索引的文档中做关键词重叠度检索，返回最相关的片段。"""
    if not _RAG_INDEX_FILE.exists():
        rag_index()
    try:
        index = json.loads(_RAG_INDEX_FILE.read_text(encoding="utf-8"))
    except Exception:  # noqa: BLE001
        return "[错误] 索引读取失败，请先调用 rag_index。"

    q_tokens = _simple_tokenize(query)
    if not q_tokens:
        return "[错误] 查询为空或无可识别关键词。"
    scores: dict[str, int] = {}
    seen: dict[str, str] = {}
    for tok in q_tokens:
        for entry in index.get(tok, []):
            scores[entry["cid"]] = scores.get(entry["cid"], 0) + 1
            seen[entry["cid"]] = entry["chunk"]
    ranked = sorted(scores.items(), key=lambda x: x[1], reverse=True)[: int(top_k)]
    if not ranked:
        return "未在文档中找到相关内容，可尝试其他关键词，或先运行 rag_index 重建索引。"
    out = []
    for cid, score in ranked:
        out.append(f"[{cid}] (匹配 {score}) {seen[cid]}")
    return "\n---\n".join(out)


# ---- 新增 Skill：跨会话长期记忆 ----
def remember(key: str, value: str) -> str:
    """把一条跨会话长期记忆写入本地知识库（key 唯一，重复会覆盖）。

    适合保存：用户偏好、项目约定、常用路径、个人信息等，
    使助手在任意新对话中都能回忆起这些信息。
    """
    if not key or not key.strip():
        return "[错误] remember 需要非空 key"
    db.upsert_memory(key.strip(), (value or "").strip())
    return f"已记住：{key.strip()} = {(value or '').strip()}"


def forget(key: str) -> str:
    """删除一条长期记忆（按 key）。"""
    if not key or not key.strip():
        return "[错误] forget 需要非空 key"
    db.delete_memory(key.strip())
    return f"已遗忘：{key.strip()}"


# ---- 新增 Skill：办公文档解析 ----
_PARSE_MAX_CHARS = 20000


def _grid_to_md(rows: list[list[str]]) -> str:
    """把二维字符串矩阵渲染为 Markdown 表格。"""
    if not rows:
        return "_(空表格)_"
    ncol = max(len(r) for r in rows)
    norm = [(r + [""] * (ncol - len(r))) for r in rows]
    md = ["| " + " | ".join(norm[0]) + " |",
          "| " + " | ".join(["---"] * ncol) + " |"]
    for r in norm[1:]:
        md.append("| " + " | ".join(r) + " |")
    return "\n".join(md)


def _parse_docx(path: str) -> str:
    from docx import Document
    from docx.document import Document as _Doc
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(path)
    lines = []
    parent = doc.element.body
    for child in parent.iterchildren():
        if child.tag == qn("w:p"):
            block = Paragraph(child, doc)
            style = (block.style.name if block.style else "") or ""
            txt = block.text.strip()
            if not txt:
                continue
            if style.startswith("Heading"):
                lvl = "".join(ch for ch in style if ch.isdigit()) or "1"
                lines.append(f"{'#' * int(lvl)} {txt}")
            elif "List" in style or "Bullet" in style:
                lines.append(f"- {txt}")
            else:
                lines.append(txt)
        elif child.tag == qn("w:tbl"):
            tbl = Table(child, doc)
            rows = [[c.text.strip().replace("\n", " ") for c in r.cells]
                    for r in tbl.rows]
            lines.append(_grid_to_md(rows))
    return "\n\n".join(lines)


def _parse_xlsx(path: str) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        out.append(f"## 工作表：{ws.title}（{ws.max_row} 行 × {ws.max_column} 列）")
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= 200:
                out.append("_(仅显示前 200 行)_")
                break
            rows.append([("" if v is None else str(v)).replace("\n", " ")
                         for v in row[:12]])
        out.append(_grid_to_md(rows) if rows else "_(空)_")
    return "\n\n".join(out)


def _parse_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"## 幻灯片 {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text.strip())
            if shape.has_table:
                tbl = shape.table
                rows = [[c.text.strip().replace("\n", " ")
                         for c in rrow.cells] for rrow in tbl.rows]
                out.append(_grid_to_md(rows))
    return "\n\n".join(out)


def _parse_pdf(path: str) -> str:
    try:
        import pymupdf as fitz
    except ImportError:  # 旧版包名仍为 fitz
        import fitz

    doc = fitz.open(path)
    pages = []
    for n, page in enumerate(doc, 1):
        pages.append(f"<!-- 第 {n} 页 -->\n{page.get_text()}")
    return "\n\n".join(pages)


def _parse_csv(path: str) -> str:
    import pandas as pd

    df = pd.read_csv(path, nrows=201)
    truncated = len(df) > 200
    md = df.head(200).to_markdown(index=False)
    if truncated:
        md += "\n\n_(仅显示前 200 行)_"
    return md


def _parse_html(path: str) -> str:
    from bs4 import BeautifulSoup

    html = Path(path).read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text("\n", strip=True)


def _parse_eml(path: str) -> str:
    import email

    with open(path, "rb") as f:
        msg = email.message_from_binary_file(f)
    parts = [f"主题: {msg.get('Subject', '')}", f"发件人: {msg.get('From', '')}",
             f"收件人: {msg.get('To', '')}", f"日期: {msg.get('Date', '')}"]
    body = ""
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                payload = part.get_payload(decode=True) or b""
                body = payload.decode(errors="replace")
                break
    else:
        payload = msg.get_payload(decode=True) or b""
        body = payload.decode(errors="replace")
    return "\n".join(parts) + "\n\n" + body


def _parse_msg(path: str) -> str:
    from extract_msg import Message

    m = Message(path)
    return (f"主题: {m.subject}\n发件人: {m.sender}\n收件人: {m.to}\n"
            f"日期: {m.date}\n\n{m.body}")


def parse_document(path: str, max_chars: int = _PARSE_MAX_CHARS) -> str:
    """解析办公/文档类文件为纯文本或 Markdown，便于模型直接阅读与检索。

    支持：.pdf .docx .xlsx .pptx .csv .html/.htm .eml .msg，
    以及纯文本类 .txt/.md/.json/.py 等（直接读取）。
    输出按 max_chars 截断并提示，避免超长内容撑爆上下文。
    """
    p = _safe_path(path)
    if not p.is_file():
        return f"[错误] 文件不存在: {path}"
    if not config.MB_SANDBOX:
        audit("parse_document", str(p))
    suffix = p.suffix.lower()
    try:
        if suffix in {".txt", ".md", ".markdown", ".json", ".py", ".log",
                      ".yaml", ".yml", ".toml"}:
            text = p.read_text(encoding="utf-8", errors="replace")
        elif suffix == ".docx":
            text = _parse_docx(str(p))
        elif suffix == ".xlsx":
            text = _parse_xlsx(str(p))
        elif suffix == ".pptx":
            text = _parse_pptx(str(p))
        elif suffix == ".pdf":
            text = _parse_pdf(str(p))
        elif suffix == ".csv":
            text = _parse_csv(str(p))
        elif suffix in {".html", ".htm"}:
            text = _parse_html(str(p))
        elif suffix == ".eml":
            text = _parse_eml(str(p))
        elif suffix == ".msg":
            text = _parse_msg(str(p))
        else:
            return (f"[不支持] 暂不支持解析 {suffix} 格式。纯文本文件可用 read_file；"
                    f"若需解析请改用 run_code 调用对应库。")
    except ImportError as e:
        return f"[错误] 缺少解析库: {e}（请用 `uv add` 安装对应依赖）"
    except Exception as e:  # noqa: BLE001
        return f"[错误] 文档解析失败: {e}"

    text = (text or "").strip()
    if not text:
        return "[内容为空] 该文档未提取到文本（可能是扫描件/图片型 PDF，需 OCR）。"
    if len(text) > int(max_chars):
        kept = int(max_chars)
        return (text[:kept]
                + f"\n\n[内容过长，已截断至 {kept} 字符（共 {len(text)} 字符）。"
                " 如需完整内容，请用 run_code 直接按页/按表读取。]")
    return text


# ---- 工具注册表 ----
TOOLS: dict[str, Callable[..., str]] = {
    "read_file": read_file,
    "write_file": write_file,
    "list_dir": list_dir,
    "search_text": search_text,
    "run_command": run_command,
    "run_code": run_code,
    "web_search": web_search,
    "manage_dir": manage_dir,
    "rag_index": rag_index,
    "rag_query": rag_query,
    "remember": remember,
    "forget": forget,
    "parse_document": parse_document,
}


async def execute_tool(name: str, arguments: dict[str, Any]) -> str:
    """异步执行工具。

    - 异步工具（如 web_search）直接 await；
    - 同步工具（run_command/run_code/文件操作等）放到线程池执行，
      避免阻塞事件循环，也让「停止」事件能在工具执行间隙被及时检测。
    """
    fn = TOOLS.get(name)
    if not fn:
        return f"[错误] 未知工具: {name}"
    try:
        if asyncio.iscoroutinefunction(fn):
            return await fn(**arguments)
        return await asyncio.to_thread(fn, **arguments)
    except TypeError as e:
        return f"[错误] 参数不匹配: {e}"
    except Exception as e:  # noqa: BLE001
        return f"[错误] 工具执行异常: {e}"


# OpenAI 兼容 tool 定义（喂给 DeepSeek）
TOOL_SCHEMAS = [
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": "读取工作区内某个文件的内容。",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "相对工作区的文件路径"}
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_file",
            "description": "在工作区内写入/创建文件。",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "相对工作区的文件路径"},
                    "content": {"type": "string", "description": "要写入的完整文本内容"},
                },
                "required": ["path", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "list_dir",
            "description": "列出工作区内某目录的内容。",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "相对路径，默认 '.' 表示工作区根"}
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "search_text",
            "description": "在工作区内按子串搜索文件内容，返回匹配的行。",
            "parameters": {
                "type": "object",
                "properties": {
                    "pattern": {"type": "string", "description": "要搜索的文本"},
                    "path": {"type": "string", "description": "搜索起点目录，默认 '.'"},
                },
                "required": ["pattern"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "run_command",
            "description": "在工作区内执行一条受控 shell 命令，返回输出。",
            "parameters": {
                "type": "object",
                "properties": {
                    "command": {"type": "string", "description": "要执行的 shell 命令"}
                },
                "required": ["command"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "run_code",
            "description": "在工作区内执行一段代码并返回结果。",
            "parameters": {
                "type": "object",
                "properties": {
                    "code": {"type": "string", "description": "代码片段"},
                    "lang": {
                        "type": "string",
                        "description": "语言，支持 python 或 js，默认 python",
                    },
                },
                "required": ["code"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "web_search",
            "description": "联网检索公开信息（无需 API key，使用 DuckDuckGo）。当用户问题需要最新/外部知识且本地无答案时调用。",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "检索关键词"},
                    "max_results": {"type": "integer", "description": "返回条数，默认 5"},
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "manage_dir",
            "description": "对指定的 TARGET_DIR 目录进行受限文件操作（增删查/移动/建目录），与通用工作区沙箱隔离。",
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {
                        "type": "string",
                        "description": "操作类型：list | read | write | mkdir | delete | move",
                    },
                    "path": {"type": "string", "description": "相对 TARGET_DIR 的路径"},
                    "content": {
                        "type": "string",
                        "description": "write 时的文件内容；move 时的目标相对路径",
                    },
                },
                "required": ["action", "path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "rag_index",
            "description": "扫描 DOC_ROOT 文档目录（md/txt/py/json/csv）建立本地检索索引。文档更新后应先调用。",
            "parameters": {"type": "object", "properties": {}},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "rag_query",
            "description": "在已索引的本地文档中检索相关片段，用于回答与用户文档相关的问题。",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "检索问题/关键词"},
                    "top_k": {"type": "integer", "description": "返回片段数，默认 5"},
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "remember",
            "description": "把一条跨会话长期记忆写入本地知识库（key 唯一，重复会覆盖）。用于保存用户偏好、项目约定、常用路径、个人信息等，使助手在任意新对话中都能回忆。",
            "parameters": {
                "type": "object",
                "properties": {
                    "key": {"type": "string", "description": "记忆键名（唯一标识，如 '用户姓名'、'项目约定'）"},
                    "value": {"type": "string", "description": "记忆内容"},
                },
                "required": ["key", "value"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "forget",
            "description": "删除一条长期记忆（按 key）。当用户明确表示不再需要某条记忆时使用。",
            "parameters": {
                "type": "object",
                "properties": {
                    "key": {"type": "string", "description": "要删除的记忆键名"},
                },
                "required": ["key"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "parse_document",
            "description": "将办公/文档类文件解析为纯文本或 Markdown，便于直接阅读与检索。支持 pdf/docx/xlsx/pptx/csv/html/eml/msg 及纯文本(.txt/.md/.json/.py 等)。大文件会按 max_chars 截断并提示。优先用本工具读取上述格式，而非 read_file 或 run_code。",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "相对工作区或绝对路径的文档路径"},
                    "max_chars": {
                        "type": "integer",
                        "description": "返回内容最大字符数，默认 20000；超过则截断",
                    },
                },
                "required": ["path"],
            },
        },
    },
]
